import zlib
import base64
import urllib.request
import os
from pptx import Presentation
from pptx.util import Inches, Pt

import time

def fetch_kroki_image(diagram_text, diagram_type, filename, retries=3):
    compressed = zlib.compress(diagram_text.encode('utf-8'), 9)
    b64 = base64.urlsafe_b64encode(compressed).decode('ascii')
    url = f"https://kroki.io/{diagram_type}/png/{b64}"
    req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
    print(f"Fetching diagram {filename}...")
    for attempt in range(retries):
        try:
            with urllib.request.urlopen(req) as response:
                with open(filename, 'wb') as f:
                    f.write(response.read())
            return
        except Exception as e:
            print(f"Attempt {attempt+1} failed: {e}. Retrying after 3s...")
            time.sleep(3)
    raise Exception(f"Failed to fetch {filename} after {retries} retries")

def add_slide(prs, title_str, content_list, is_title_slide=False):
    if is_title_slide:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_str
        subtitle.text = content_list[0]
    else:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_str
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
    return slide

def add_image_slide(prs, title_str, image_path):
    slide_layout = prs.slide_layouts[5] # Title Only slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_str
    
    # We will center the image horizontally by calculating its display width
    # 8 inches width, leaving 1 inch on each side for a 10 inch slide width
    slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
    return slide

# Define Diagram Texts
usecase_text = """@startuml
skinparam packageStyle rectangle
actor Employee
actor Manager
actor Admin

package "RevWorkforce" {
  usecase "Apply Leave" as UC1
  usecase "Self Review" as UC2
  usecase "View Goals" as UC3
  
  usecase "Approve Leaves" as UC4
  usecase "Review Direct Reports" as UC5
  usecase "Assign Goals" as UC6
  
  usecase "Settings & Reports" as UC7
}

Employee --> UC1
Employee --> UC2
Employee --> UC3

Manager --> UC1
Manager --> UC2
Manager --> UC3
Manager --> UC4
Manager --> UC5
Manager --> UC6

Admin --> UC1
Admin --> UC2
Admin --> UC3
Admin --> UC4
Admin --> UC5
Admin --> UC6
Admin --> UC7
@enduml
"""

erd_text = """erDiagram
    EMPLOYEE {
        Long employeeId PK
        String name
        Long managerId FK
    }
    LEAVE_BALANCE {
        Long id PK
        Long employeeId FK
        Int casualLeave
    }
    LEAVE_REQUEST {
        Long id PK
        Long employeeId FK
        Date fromDate
        String status
    }
    PERFORMANCE_REVIEW {
        Long reviewId PK
        Long employeeId FK
        String status
    }
    GOAL {
        Long goalId PK
        Long reviewId FK
        String description
    }

    EMPLOYEE ||--|| LEAVE_BALANCE : "has"
    EMPLOYEE ||--o{ LEAVE_REQUEST : "submits"
    EMPLOYEE ||--o{ PERFORMANCE_REVIEW : "requires"
    EMPLOYEE ||--o{ GOAL : "assigned"
    PERFORMANCE_REVIEW |o--o{ GOAL : "links to"
    EMPLOYEE ||--o{ EMPLOYEE : "reports to"
"""

class_text = """classDiagram
    class LeaveController {
        +applyLeave(LeaveRequest request)
        +getMyLeaves() List
    }
    class LeaveService {
        +applyLeave(Long empId, LeaveRequest request)
    }
    class LeaveRepository {
        +findByEmployee(Long empId) List
    }
    class LeaveBalance {
        -Long id
        -Integer casualLeave
    }

    LeaveController --> LeaveService : "injects & uses"
    LeaveService --> LeaveRepository : "injects & uses"
    LeaveRepository --> LeaveBalance : "queries"
"""

sequence_text = """sequenceDiagram
    actor Employee
    participant Controller as LeaveController
    participant Service as LeaveService
    participant Repo as LeaveRepository
    participant DB as Database

    Employee->>Controller: POST /api/leaves/apply
    Controller->>Service: applyLeave(empId, request)
    Service->>Repo: checkBalance(empId)
    Repo->>DB: Query LeaveBalance
    DB-->>Repo: Balance Data
    Repo-->>Service: LeaveBalance
    Service->>Service: Validate enough days left
    Service->>Repo: save(LeaveApplication)
    Repo-->>Service: Saved Request
    Service-->>Controller: DTO Formatted
    Controller-->>Employee: 201 Created (Success)
"""

# Fetch diagram images
# Use provided image for Use Case: C:/Users/nihar/.gemini/antigravity/brain/635fbe12-3cd2-4b08-afc9-8aaba187d435/uploaded_media_1772609169667.png
fetch_kroki_image(erd_text, "mermaid", "erd.png")
fetch_kroki_image(class_text, "mermaid", "class.png")
fetch_kroki_image(sequence_text, "mermaid", "sequence.png")

# Create Presentation
prs = Presentation()

add_slide(prs, "RevWorkforce: Empowering Modern Teams", ["A Comprehensive HR & Performance Management Solution"], is_title_slide=True)

add_slide(prs, "Introduction to RevWorkforce", [
    "A simple, all-in-one HR platform designed to make managing people easier.",
    "Empowers employees to request time off, track goals, and review their performance.",
    "Provides managers with tools to approve leaves and guide their team's growth.",
    "Gives admins control over system settings and company-wide reports."
])

add_slide(prs, "Our Technology Stack", [
    "Backend Engine: Powered by Java & Spring Boot to handle business rules securely.",
    "Database: Uses MySQL via Spring Data JPA to store all HR data safely.",
    "Security: Uses stateless JWT tokens to ensure only authorized people access the data.",
    "Web Interface: Built with clean, simple frontend templates for a great user experience."
])

add_slide(prs, "Application Layers Explained", [
    "1. Controller Layer: The 'Front Desk'. It receives user requests from the web and sends back responses.",
    "2. Service Layer: The 'Brain'. It contains all the business rules (e.g., checking if someone has enough leave balance).",
    "3. Repository Layer: The 'Filing Cabinet'. It talks directly to the database to save and retrieve information."
])

# Insert diagram slides
add_image_slide(prs, "Use Case Diagram", r"C:\Users\nihar\.gemini\antigravity\brain\635fbe12-3cd2-4b08-afc9-8aaba187d435\media__1772616600676.png")
add_image_slide(prs, "Entity Relationship Diagram (ERD)", "erd.png")
add_image_slide(prs, "Architecture Class Diagram", "class.png")
add_image_slide(prs, "Sequence Diagram (Leave Request)", "sequence.png")

# Remaining slides
add_slide(prs, "The Employee Experience", [
    "Leave Management: Easily apply for leaves and automatically track balances dynamically.",
    "Performance Reviews: Submit comprehensive self-assessments and track personal accomplishments.",
    "Goal Tracking: Keep track of personal and assigned goals directly linked to appraisal cycles."
])

add_slide(prs, "Powerful Manager Capabilities", [
    "Team Leadership: Gain complete visibility into team structure and leave calendars.",
    "Efficient Approvals: Seamlessly approve or reject leave requests.",
    "Performance Management: Access direct reports' reviews, provide feedback, and assign ratings."
])

add_slide(prs, "Robust Admin Control", [
    "System Configuration: Manage critical application data and organization-wide announcements.",
    "Powerful Reporting: Generate and export Employee summaries and Leave Utilization reports.",
    "User Management: Maintain robust employee onboarding configurations and access controls."
])

add_slide(prs, "Security & Compliance", [
    "Stateless Authentication: Secure API endpoints to protect sensitive HR and performance data.",
    "Audit Logging: Automatic tracking and recording of critical system actions for organizational compliance.",
    "Strict Access Control: Rigid separation of concerns ensuring that employees only see their data."
])

add_slide(prs, "Looking Ahead (Roadmap)", [
    "Advanced analytics and AI-driven insights for HR strategy.",
    "Deep integration with external calendar providers (Google Workspace / Microsoft 365).",
    "Dedicated Mobile Application for on-the-go HR management access."
])

add_slide(prs, "Thank You!", [
    "Q&A Session",
    "Live Demonstration"
])

prs.save('RevWorkforce_Presentation_v6.pptx')
print("Successfully generated diagram-enhanced RevWorkforce_Presentation_v6.pptx")
